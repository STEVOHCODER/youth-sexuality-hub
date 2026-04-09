# filename: generate_rwanda_mining_ppt.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_rwanda_mining_ppt():
    """
    Generates a PowerPoint presentation on the Critical Review of Mining Laws in Rwanda.
    Note: python-pptx does not natively support adding slide animations or transitions via Python.
    The slides are structured cleanly so that you can simply press Ctrl+A and apply a standard 
    animation (like 'Fade' or 'Fly In') to all bullet points in PowerPoint once generated.
    """
    # Initialize presentation
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define a professional, earth-toned color palette suitable for mining/governance
    TITLE_COLOR = RGBColor(46, 64, 83)     # Dark Slate (Professional/Legal)
    TEXT_COLOR = RGBColor(64, 64, 64)      # Dark Gray (Readability)
    ACCENT_COLOR = RGBColor(212, 172, 13)  # Goldenrod/Mineral Accent
    
    # Slide Layouts
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    
    # Helper function to style titles
    def style_title(shape, text, is_main_title=False):
        shape.text = text
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if is_main_title else PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.bold = True
                run.font.name = 'Calibri'
                run.font.size = Pt(44) if is_main_title else Pt(36)
                run.font.color.rgb = TITLE_COLOR

    # Helper function to style body/bullet points
    def style_body(shape, points):
        text_frame = shape.text_frame
        text_frame.clear()  # Clear default text
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.space_after = Pt(14)
            p.level = 0
            for run in p.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(22)
                run.font.color.rgb = TEXT_COLOR

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    style_title(title, "A Critical Review of Mining Laws in Rwanda", is_main_title=True)
    
    subtitle.text = "Evaluating Regulatory Frameworks, Challenges, and Sustainable Practices\n\nGenerated for Policy & Legal Analysis"
    for p in subtitle.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(20)
            run.font.color.rgb = ACCENT_COLOR

    # --- Slide Content Data ---
    slides_data = [
        {
            "title": "1. Introduction to Rwanda's Mining Sector",
            "bullets": [
                "Mining is Rwanda’s second-largest export revenue earner, crucial for economic growth.",
                "Primary minerals include the '3Ts': Tin (Cassiterite), Tungsten (Wolfram), and Tantalum (Coltan), alongside gold and gemstones.",
                "The sector has transitioned from a largely informal, artisanal state to a heavily regulated, professionalizing industry.",
                "Legal frameworks aim to balance aggressive economic targets with sustainable environmental practices."
            ]
        },
        {
            "title": "2. Evolution of the Legal Framework",
            "bullets": [
                "Post-colonial laws heavily favored state control with minimal environmental oversight.",
                "2014 Mining Law: Introduced sweeping reforms focusing on privatization, licensing, and traceability.",
                "2018 Revised Law: Stricter regulations on value addition, environmental rehabilitation, and formalization of Artisanal and Small-scale Mining (ASM).",
                "Establishment of the Rwanda Mines, Petroleum and Gas Board (RMB) in 2017 to centralize regulation and monitoring."
            ]
        },
        {
            "title": "3. Strengths of the Current Mining Laws",
            "bullets": [
                "Traceability & Transparency: Mandatory tagging (e.g., iTSCi scheme) ensures conflict-free mineral exports.",
                "Value Addition: The law heavily incentivizes domestic processing and smelting to maximize export value.",
                "Environmental Mandates: Strict requirement for Environmental Impact Assessments (EIAs) and rehabilitation bonds.",
                "Decentralization: Licensing frameworks structured to include local government oversight."
            ]
        },
        {
            "title": "4. The Challenge of ASM Formalization",
            "bullets": [
                "Artisanal and Small-scale Miners (ASM) produce a massive portion of the country's minerals but struggle with compliance.",
                "Bureaucratic Licensing: The process for obtaining legal mining rights is often too costly and complex for local artisanal miners.",
                "Access to Finance: Strict legal requirements are difficult to meet without capital, yet banks hesitate to fund ASM operations.",
                "Result: A persistence of illegal mining operations despite strong regulatory texts."
            ]
        },
        {
            "title": "5. Environmental & Social Disconnects",
            "bullets": [
                "While the law mandates rehabilitation bonds, enforcement and post-mining land restoration often fall short.",
                "Land Use Conflicts: Tension between agricultural land needs and mining concessions in a densely populated country.",
                "Occupational Health & Safety (OHS): The law sets high safety standards, but frequent mine collapses indicate enforcement gaps.",
                "Corporate Social Responsibility (CSR): Community benefit-sharing mechanisms are legally ambiguous and inconsistently applied."
            ]
        },
        {
            "title": "6. Institutional & Capacity Constraints",
            "bullets": [
                "The Rwanda Mines, Petroleum and Gas Board (RMB) has strong statutory powers but faces capacity limits in field inspections.",
                "Geological Data: Insufficient state-funded geological surveys leave exploration risks entirely to private investors.",
                "Judicial & Dispute Resolution: Lack of specialized mining courts can delay the resolution of concession disputes.",
                "Taxation & Royalties: High compliance costs can disincentivize full declaration of yields by smaller mining entities."
            ]
        },
        {
            "title": "7. Recommendations for Legal Reform",
            "bullets": [
                "Streamlined ASM Licensing: Create a tiered, simplified licensing system specifically tailored to empower artisanal miners.",
                "Financial Support Mechanisms: Establish state-backed guarantee funds to help local miners acquire modern equipment.",
                "Stricter Enforcement with Capacity Building: Shift from purely punitive measures to compliance-assistance programs.",
                "Clarify CSR Mandates: Legislate specific, measurable percentages of revenue to be reinvested directly into local communities."
            ]
        },
        {
            "title": "8. Conclusion",
            "bullets": [
                "Rwanda's mining laws are highly progressive on paper, pioneering conflict-free traceability and value addition.",
                "The critical gap lies in the implementation and the socio-economic realities of the artisanal sector.",
                "Future legal iterations must focus on inclusivity, ensuring the regulatory burden does not marginalize local miners.",
                "Bridging the gap between the law and on-the-ground realities will unlock the sector's true sustainable potential."
            ]
        }
    ]

    # --- Generate Content Slides ---
    for slide_info in slides_data:
        slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        style_title(title_shape, slide_info["title"])
        style_body(body_shape, slide_info["bullets"])

    # Save the resulting presentation
    output_filename = "Rwanda_Mining_Laws_Review.pptx"
    prs.save(output_filename)
    print(f"Presentation generated successfully: {output_filename}")
    print("Note: To add animations, open the PPTX file, go to 'Slide Sorter' view, select all slides, go to 'Animations', and select 'Fade' or 'Appear' for the bullet points.")

if __name__ == "__main__":
    create_rwanda_mining_ppt()
