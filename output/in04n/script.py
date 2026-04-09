import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Define Corporate Colors
NAVY_BLUE = RGBColor(31, 73, 125)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)

def add_header_accent(slide, prs):
    """Adds a consistent corporate navy accent bar to the top of slides."""
    left = top = Inches(0)
    width = prs.slide_width
    height = Inches(0.15)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = NAVY_BLUE
    accent.line.fill.background() # No border

def format_title(title_shape):
    """Standardizes slide titles to Bold Navy Calibri."""
    if title_shape.has_text_frame:
        for p in title_shape.text_frame.paragraphs:
            p.font.name = 'Calibri'
            p.font.bold = True
            p.font.color.rgb = NAVY_BLUE

def generate_professional_ppt(file_name="Corporate_Financial_Deck.pptx"):
    prs = Presentation()

    # ==========================================
    # SLIDE 1: Title Slide (Layout 0)
    # ==========================================
    slide_layout_title = prs.slide_layouts[0]
    slide_1 = prs.slides.add_slide(slide_layout_title)
    
    title = slide_1.shapes.title
    subtitle = slide_1.placeholders[1]
    
    title.text = "Q1 Corporate Financial Report"
    subtitle.text = "Regional Revenue & Margin Analysis\nGenerated Automatically"
    
    format_title(title)
    # Format subtitle
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = DARK_GRAY
        p.font.size = Pt(18)

    # Add decorative geometric shapes to Title Slide
    shape = slide_1.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(1), Inches(3.5), Inches(0.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_BLUE
    shape.line.fill.background()

    # ==========================================
    # SLIDE 2: Executive Summary (Layout 1)
    # ==========================================
    slide_layout_content = prs.slide_layouts[1]
    slide_2 = prs.slides.add_slide(slide_layout_content)
    add_header_accent(slide_2, prs)
    
    slide_2.shapes.title.text = "Executive Summary"
    format_title(slide_2.shapes.title)
    
    tf = slide_2.placeholders[1].text_frame
    tf.text = "Strong performance across all four global regions."
    
    # Add bullet hierarchy
    p1 = tf.add_paragraph()
    p1.text = "North America led the Enterprise License segment."
    p1.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "EMEA showed a 15% increase in Cloud Storage attach rates."
    p2.level = 1
    
    p3 = tf.add_paragraph()
    p3.text = "Profit margins stabilized at ~45%."
    p3.level = 0
    
    p4 = tf.add_paragraph()
    p4.text = "Hardware Setup costs reduced due to supply chain optimizations."
    p4.level = 1

    # Standardize bullet fonts
    for p in tf.paragraphs:
        p.font.name = 'Calibri'
        p.font.size = Pt(20 if p.level == 0 else 18)
        p.font.color.rgb = DARK_GRAY

    # ==========================================
    # SLIDE 3: Financial Highlights (Layout 3 - Two Column)
    # ==========================================
    slide_layout_two_col = prs.slide_layouts[3]
    slide_3 = prs.slides.add_slide(slide_layout_two_col)
    add_header_accent(slide_3, prs)
    
    slide_3.shapes.title.text = "Financial Highlights"
    format_title(slide_3.shapes.title)
    
    # Left Column (Revenue)
    left_tf = slide_3.placeholders[1].text_frame
    left_tf.text = "Revenue Metrics"
    left_tf.paragraphs[0].font.bold = True
    
    p = left_tf.add_paragraph()
    p.text = "• Gross Revenue exceeded $45M"
    p.level = 1
    p = left_tf.add_paragraph()
    p.text = "• Consulting Services drove 30% of growth"
    p.level = 1
    
    # Right Column (Profit)
    right_tf = slide_3.placeholders[2].text_frame
    right_tf.text = "Profitability & COGS"
    right_tf.paragraphs[0].font.bold = True
    
    p = right_tf.add_paragraph()
    p.text = "• Net Profit margins expanded"
    p.level = 1
    p = right_tf.add_paragraph()
    p.text = "• COGS variance restricted to <5%"
    p.level = 1

    for tf in [left_tf, right_tf]:
        for p in tf.paragraphs:
            p.font.name = 'Calibri'
            p.font.color.rgb = DARK_GRAY

    # ==========================================
    # SLIDE 4: Regional Summary Data (Table)
    # ==========================================
    slide_layout_blank = prs.slide_layouts[5] # Title only
    slide_4 = prs.slides.add_slide(slide_layout_blank)
    add_header_accent(slide_4, prs)
    
    slide_4.shapes.title.text = "Regional Profitability Summary"
    format_title(slide_4.shapes.title)
    
    # Create Table
    rows, cols = 5, 4
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(2)
    table_shape = slide_4.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set Column Widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)
    
    headers = ["Region", "Revenue ($)", "Profit ($)", "Margin (%)"]
    data = [
        ["North America", "12,450,000", "5,602,500", "45.0%"],
        ["EMEA", "9,800,000", "4,116,000", "42.0%"],
        ["APAC", "14,200,000", "6,816,000", "48.0%"],
        ["LATAM", "5,100,000", "2,193,000", "43.0%"]
    ]
    
    # Format Headers
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
            
    # Format Data Rows with slight alternating visual (Zebra)
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = DARK_GRAY
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

    # ==========================================
    # SLIDE 5: Next Steps (Process Arrows)
    # ==========================================
    slide_5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_header_accent(slide_5, prs)
    
    slide_5.shapes.title.text = "Strategic Roadmap"
    format_title(slide_5.shapes.title)

    # Process Arrows (Chevrons)
    steps = ["Q1 Review", "Strategy Pivot", "Q2 Execution"]
    arrow_width = Inches(2.5)
    arrow_height = Inches(1)
    top_pos = Inches(3)
    
    for i, step_text in enumerate(steps):
        left_pos = Inches(1 + (i * 2.8)) # Space them out
        shape = slide_5.shapes.add_shape(MSO_SHAPE.CHEVRON, left_pos, top_pos, arrow_width, arrow_height)
        
        # Style Arrow
        shape.fill.solid()
        # Alternate shades of Navy for process flow
        intensity = max(31 + (i * 40), 0)
        shape.fill.fore_color.rgb = RGBColor(intensity, 73, 125)
        shape.line.fill.background()
        
        # Add Text to Arrow
        tf = shape.text_frame
        tf.text = step_text
        for p in tf.paragraphs:
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER

    prs.save(file_name)
    print(f"Successfully generated PowerPoint deck: '{file_name}'")

if __name__ == "__main__":
    generate_professional_ppt()
