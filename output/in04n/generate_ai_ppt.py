# filename: generate_ai_ppt.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ai_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define a clean, modern color palette
    TITLE_COLOR = RGBColor(0, 51, 102)     # Deep Corporate Blue
    TEXT_COLOR = RGBColor(64, 64, 64)      # Dark Grey for readability
    ACCENT_COLOR = RGBColor(0, 120, 215)   # Vibrant Light Blue
    
    # Slide Layouts
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    BULLET_SLIDE_LAYOUT = prs.slide_layouts[1]
    
    # Helper function to style titles
    def style_title(shape, text, is_main_title=False):
        shape.text = text
        text_frame = shape.text_frame
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER if is_main_title else PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.bold = True
                run.font.name = 'Arial'
                run.font.size = Pt(48) if is_main_title else Pt(36)
                run.font.color.rgb = TITLE_COLOR

    # Helper function to style body/bullet points
    def style_body(shape, points):
        text_frame = shape.text_frame
        text_frame.clear()  # Clear default text
        
        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.space_after = Pt(18)
            p.level = 0
            for run in p.runs:
                run.font.name = 'Arial'
                run.font.size = Pt(24)
                run.font.color.rgb = TEXT_COLOR

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    style_title(title, "The Advantages of Artificial Intelligence", is_main_title=True)
    
    subtitle.text = "Unlocking Potential, Driving Innovation, and Transforming the Future"
    for p in subtitle.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = 'Arial'
            run.font.size = Pt(24)
            run.font.color.rgb = ACCENT_COLOR

    # --- Slide Content Data ---
    slides_data = [
        {
            "title": "1. What is AI?",
            "bullets": [
                "Machines designed to perform tasks that typically require human intelligence.",
                "Encompasses Machine Learning (ML), Deep Learning, and Natural Language Processing (NLP).",
                "Represents a fundamental shift in computing, moving from rule-based to learning-based systems.",
                "Seamlessly integrates into daily operations to enhance human capabilities."
            ]
        },
        {
            "title": "2. Unmatched Automation & Efficiency",
            "bullets": [
                "Automates repetitive and mundane tasks, freeing human workers for creative and strategic roles.",
                "Operates 24/7 without fatigue, breaks, or degradation in quality.",
                "Drastically reduces operational bottlenecks and minimizes human error.",
                "Streamlines complex supply chains and manufacturing processes."
            ]
        },
        {
            "title": "3. Enhanced Decision Making",
            "bullets": [
                "Processes and analyzes vast amounts of Big Data in mere seconds.",
                "Identifies hidden patterns, correlations, and market trends that humans might miss.",
                "Provides accurate predictive analytics for proactive business maneuvers.",
                "Removes emotional bias from critical data interpretation and forecasting."
            ]
        },
        {
            "title": "4. Hyper-Personalization",
            "bullets": [
                "Analyzes consumer behavior to deliver highly tailored product and content recommendations.",
                "Dramatically enhances customer satisfaction, engagement, and brand loyalty.",
                "Powers intelligent virtual assistants and chatbots for instantaneous 24/7 customized support.",
                "Optimizes marketing campaigns to ensure maximum conversion rates and ROI."
            ]
        },
        {
            "title": "5. Revolutionizing Healthcare",
            "bullets": [
                "Accelerates drug discovery, reducing research and development time from years to months.",
                "Improves diagnostic accuracy through advanced medical image analysis (e.g., X-rays, MRIs).",
                "Enables precision medicine tailored to individual patient genetic profiles.",
                "Facilitates continuous remote patient monitoring via smart wearables and IoT devices."
            ]
        },
        {
            "title": "6. Driving Innovation & Economic Growth",
            "bullets": [
                "Acts as a powerful catalyst for entirely new industries and disruptive business models.",
                "Enhances Research & Development (R&D) across engineering, software, and physical sciences.",
                "Optimizes global energy consumption and heavily contributes to sustainability efforts.",
                "Projected to add trillions of dollars to the global economy over the next decade."
            ]
        },
        {
            "title": "Conclusion",
            "bullets": [
                "AI is a powerful tool for human empowerment, augmenting rather than replacing us.",
                "The key to unlocking its full potential lies in ethical and responsible implementation.",
                "Continuous learning and adaptation are essential for workforces to thrive in an AI-driven world.",
                "Embracing AI today effectively prepares organizations for the complex challenges of tomorrow."
            ]
        }
    ]

    # --- Generate Content Slides ---
    for slide_info in slides_data:
        slide = prs.slides.add_slide(BULLET_SLIDE_LAYOUT)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        style_title(title_shape, slide_info["title"])
        style_body(body_shape, slide_info["bullets"])

    # Save the resulting presentation
    output_filename = "Advantages_of_AI_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation generated successfully: {output_filename}")

if __name__ == "__main__":
    create_ai_presentation()
