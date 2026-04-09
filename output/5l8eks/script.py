from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # ---------------------------------------------------
    # Slide 1: Title Slide
    # ---------------------------------------------------
    slide_1 = prs.slides.add_slide(title_slide_layout)
    title_1 = slide_1.shapes.title
    subtitle_1 = slide_1.placeholders[1]
    
    title_1.text = "Advanced Mathematical Modeling & Scientific Architecture"
    subtitle_1.text = "Universal Architect \nFractional Calculus, Statistical Physics & Data Automation"
    
    # ---------------------------------------------------
    # Slide 2: Core Capabilities
    # ---------------------------------------------------
    slide_2 = prs.slides.add_slide(bullet_slide_layout)
    title_2 = slide_2.shapes.title
    title_2.text = "Core Architectural Disciplines"
    
    body_2 = slide_2.shapes.placeholders[1]
    tf_2 = body_2.text_frame
    
    p = tf_2.text = "1. Mathematical Modeling"
    p = tf_2.add_paragraph()
    p.text = "Expertise in Fractional Calculus, Riemann-Liouville, and Caputo derivatives."
    p.level = 1
    
    p = tf_2.add_paragraph()
    p.text = "2. Full-Stack Development"
    p.level = 0
    p = tf_2.add_paragraph()
    p.text = "Production-ready architectures using React, FastAPI, and robust databases."
    p.level = 1

    p = tf_2.add_paragraph()
    p.text = "3. Data & Excel Automation"
    p.level = 0
    p = tf_2.add_paragraph()
    p.text = "Complex pivot tables, dashboarding, and analytics via Pandas and Openpyxl."
    p.level = 1
    
    # ---------------------------------------------------
    # Slide 3: Fractional Calculus in Physics
    # ---------------------------------------------------
    slide_3 = prs.slides.add_slide(bullet_slide_layout)
    title_3 = slide_3.shapes.title
    title_3.text = "Fractional Calculus in Statistical Physics"
    
    body_3 = slide_3.shapes.placeholders[1]
    tf_3 = body_3.text_frame
    
    tf_3.text = "Why use Fractional Calculus?"
    p = tf_3.add_paragraph()
    p.text = "Memory Effects: Unlike integer-order derivatives, fractional derivatives depend on the history of the function."
    p.level = 1
    
    p = tf_3.add_paragraph()
    p.text = "Anomalous Diffusion:"
    p.level = 0
    p = tf_3.add_paragraph()
    p.text = "Standard diffusion models fail in complex media (e.g., biological tissues, porous materials)."
    p.level = 1
    p = tf_3.add_paragraph()
    p.text = "Fractional differential equations accurately model sub-diffusion and super-diffusion."
    p.level = 1

    # ---------------------------------------------------
    # Slide 4: Next Steps & Customization
    # ---------------------------------------------------
    slide_4 = prs.slides.add_slide(bullet_slide_layout)
    title_4 = slide_4.shapes.title
    title_4.text = "Tailoring to Your Requirements"
    
    body_4 = slide_4.shapes.placeholders[1]
    tf_4 = body_4.text_frame
    
    tf_4.text = "How to proceed:"
    
    p = tf_4.add_paragraph()
    p.text = "Provide a Specific Topic: Let me know the exact scientific or business focus."
    p.level = 1
    
    p = tf_4.add_paragraph()
    p.text = "Attach Datasets: I can ingest CSV/Excel files directly, analyze them, and build dynamic slides with charts based on your data."
    p.level = 1
    
    p = tf_4.add_paragraph()
    p.text = "Define Visuals: Request specific plots (using matplotlib/scipy) to be embedded directly into the slides."
    p.level = 1

    # Save the presentation
    output_filename = "Scientific_Modeling_Overview.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully generated and saved as '{output_filename}'")

if __name__ == "__main__":
    create_presentation()
